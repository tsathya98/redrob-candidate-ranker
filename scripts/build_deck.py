#!/usr/bin/env python3
"""Fill the organizer-provided idea-submission template with our content (preserves their design).

We fill the body text box of each slide (and the title-page fields) in place, keeping the template's
layout/branding. Content mirrors docs/09 + README. Re-run after editing CONTENT below.

  uv pip install python-pptx
  uv run python scripts/build_deck.py
  -> submission/Redrob_Idea_Submission.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


def _pPr(p, marL=None, indent=None):
    pPr = p._p.get_or_add_pPr()
    if marL is not None:
        pPr.set("marL", str(Inches(marL)))
    if indent is not None:
        pPr.set("indent", str(Inches(indent)))
    for t in ("a:buClr", "a:buFont", "a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(t)):
            pPr.remove(e)
    return pPr


def no_bullet(p, marL=0.0):
    """Flush-left paragraph, no bullet (headers / notes)."""
    pPr = _pPr(p, marL=marL, indent=0.0)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def bullet(p, char, marL, indent, color="3730A3"):
    """Clean hanging bullet: marker at the indent, text wraps aligned at marL."""
    pPr = _pPr(p, marL=marL, indent=indent)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    buClr.append(buClr.makeelement(qn("a:srgbClr"), {"val": color}))
    pPr.append(buClr)
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))

TEMPLATE = Path("submission/Idea_Submission_Template_Redrob.pptx")
OUT = Path("submission/Redrob_Idea_Submission.pptx")

BODY = RGBColor(0x1F, 0x29, 0x37)   # ink / body text
HEAD = RGBColor(0x37, 0x30, 0xA3)   # indigo accent (matches the INDIA.RUNS brand)
MUTE = RGBColor(0x6B, 0x72, 0x80)
PANEL = RGBColor(0xF6, 0xF7, 0xFB)  # light content-card fill
BORDER = RGBColor(0xE2, 0xE6, 0xF0)  # card border
CARD = RGBColor(0xFF, 0xFF, 0xFF)   # stat-card fill
GREEN = RGBColor(0x16, 0xA3, 0x4A)  # "live" accent

# Per-slide content. Each item: ("h"|"b"|"s"|"note", text)
#   h = bold sub-header, b = bullet, s = sub-bullet (indented), note = muted italic
CONTENT = {
    2: [
        ("h", "What is your proposed solution?"),
        ("b", "A hybrid, explainable, CPU-only ranking engine: honeypot filter -> JD gated rubric -> "
              "semantic + lexical relevance -> cross-encoder rerank -> behavioral multiplier -> "
              "fact-grounded reasoning."),
        ("h", "What makes it different from keyword matching?"),
        ("b", "Substance over keywords: scores what candidates BUILT (career history); the skills list is "
              "excluded, so keyword-stuffers don't rank."),
        ("b", "Reads profiles: an impossibility/honeypot filter removes logically-impossible profiles "
              "(0 in our top-100 vs the >10% that disqualifies)."),
        ("b", "Availability-aware: a bounded behavioral multiplier down-weights unreachable "
              "'perfect-on-paper' candidates."),
        ("b", "Compliant by design: heavy models run offline; the scored step is stdlib, CPU-only, "
              "offline, ~150s."),
    ],
    3: [
        ("h", "Key requirements extracted from the JD"),
        ("b", "Must-have: production embeddings retrieval, vector/hybrid search, strong Python, ranking "
              "evaluation (NDCG/MRR/MAP/A-B)."),
        ("b", "Ideal profile: 5-9 yrs, product company (not pure services/research), shipped an end-to-end "
              "ranking/search/recsys system, Pune/Noida."),
        ("b", "Negatives: services-only, CV/speech-without-NLP, title-chasers, recent-LLM-only."),
        ("h", "Evaluating fit beyond keyword matching"),
        ("b", "Career-evidence of retrieval/ranking work + an engineering-role gate ('Marketing Manager' "
              "+ AI skills ~ 0)."),
        ("b", "Skill DEPTH (proficiency x endorsements/duration x assessments), not count; embeddings "
              "catch plain-language Tier-5s; 23 Redrob signals weigh availability."),
    ],
    4: [
        ("h", "Retrieve -> score -> rerank"),
        ("b", "Recall funnel: a fast structured score over all 100k keeps a top pool."),
        ("b", "Semantic: bge-small-en-v1.5 cosine to a distilled JD query, blended 50/50 with lexical "
              "concept matching."),
        ("b", "Rerank: bge-reranker-v2-m3 cross-encoder takes half-authority over the top tier "
              "(drives NDCG@10)."),
        ("b", "Weighted components: relevance .30, title/career .28, skill-depth .20, experience .12, "
              "education .10."),
        ("b", "Then multiply by soft penalties (services/CV-speech/title-chaser/location) and a bounded "
              "behavioral availability factor (0.5 to 1.15)."),
        ("b", "Final: score = fit x penalties x behavioral; tie-break by candidate_id. Models run offline; "
              "rank.py reads cached JSON."),
    ],
    5: [
        ("h", "How decisions are explained"),
        ("b", "Each candidate shows a component breakdown + relevance parts (lexical / semantic / "
              "cross-encoder) and a 1-2 sentence reasoning."),
        ("h", "No hallucinations"),
        ("b", "Reasoning is templated strictly from real profile facts (no LLM), so it cannot cite anything "
              "not in the profile; tone matches rank."),
        ("h", "Suspicious profiles"),
        ("b", "Honeypot filter (e.g. 'expert' skill with 0 months used). 65 flagged, 0 false positives, "
              "0 reach the top-100."),
    ],
    6: [
        ("h", "Complete workflow: JD input -> ranked output"),
        ("b", "OFFLINE (GPU ok, one-time): embed the candidate pool + cross-encoder rerank -> cache "
              "artifacts/*.json."),
        ("b", "SCORED (rank.py, CPU/offline, <=5 min): stream 100k -> honeypot filter -> component scoring "
              "x behavioral multiplier -> top-100 heap -> fact-grounded reasoning -> submission.csv."),
        ("b", "Validate: data/validate_submission.py confirms format before upload."),
        ("note", "Full workflow diagram: docs/05_approach_and_roadmap.md"),
    ],
    # 7 (System Architecture) is drawn as a native diagram - see draw_arch().
    8: [
        ("h", "Results demonstrating ranking quality"),
        ("b", "Validator passes (100 rows, ranks 1-100, scores non-increasing, tie-break by id)."),
        ("b", "Honeypots in top-100: 0 - clears the >10% hard disqualifier."),
        ("b", "Quality (proxy, ground truth is hidden): top-100 -> 83 in-band, 97 show ranking work, only "
              "1 junior, 92 available. Top-10 are senior product-company retrieval/ranking/NLP engineers "
              "(Meta, Apple, Netflix, Zomato, Paytm, Razorpay...), 5.9-7.9 yrs."),
        ("h", "Meeting the runtime / compute constraints"),
        ("b", "rank.py: ~150s, CPU-only, offline, <16 GB (the only GPU step is offline precompute)."),
        ("note", "Method validated by a 3-way query ablation + full-100k negative calibration (docs/09)."),
    ],
    9: [
        ("h", "Technologies & why"),
        ("b", "Ranker: Python 3.13, stdlib-only - deliberate, guarantees CPU/offline reproducibility at Stage 3."),
        ("b", "Offline precompute: PyTorch + sentence-transformers; bge-small-en-v1.5 (embeddings) + "
              "bge-reranker-v2-m3 (cross-encoder) - open-source, no hosted API (Cohere excluded: network)."),
        ("b", "Demo: FastAPI; Vite + React + TypeScript + Tailwind; Recharts."),
        ("b", "Tooling: uv (reproducible envs, no pip); just (cross-platform tasks); Docker (HF Space)."),
        ("note", "Every choice keeps the scored path inside the compute budget while heavy models run offline."),
    ],
    10: [
        ("h", "Submission assets"),
        ("b", "GitHub (public): full code + docs; reproduce with just rank."),
        ("b", "Ranked output: submission.xlsx (top-100)."),
        ("b", "This deck (PDF): approach, methodology, results."),
        ("b", "README + docs/00-09 + PROJECT_LOG (iteration trail)."),
        ("b", "Live sandbox: hosted Hugging Face Docker Space (also: just serve / docker-run)."),
    ],
    # 11 (closing) - left as the template's clean "THANK YOU" branded slide (don't overwrite).
}

TITLE_FIELDS = {  # slide 1: prefix -> filled value
    "Team Name": "Argmax",
    "Team Leader Name": "Sathya T",
    "Problem Statement": ("Rank the top-100 best-fit candidates from a 100,000-profile pool for the "
                          "Senior AI Engineer (Founding Team) @ Redrob AI role - contextual + behavioral "
                          "fit over keyword overlap, CPU-only and offline (<=5 min)."),
}


def style(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def fill_body(tf, items):
    tf.word_wrap = True
    tf.clear()
    first = True
    for kind, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.02
        if kind == "h":
            no_bullet(p, marL=0.0)
            p.space_before = Pt(9); p.space_after = Pt(4)
            r = p.add_run(); r.text = text; style(r, 13.5, HEAD, bold=True)
        elif kind == "note":
            no_bullet(p, marL=0.0)
            p.space_before = Pt(4); p.space_after = Pt(2)
            r = p.add_run(); r.text = text; style(r, 10, MUTE, italic=True)
        elif kind == "s":
            bullet(p, "-", 0.55, -0.27, color="6B7280")
            p.space_after = Pt(3)
            r = p.add_run(); r.text = text; style(r, 10.5, BODY)
        else:  # bullet
            bullet(p, "•", 0.28, -0.28)
            p.space_after = Pt(5)
            r = p.add_run(); r.text = text; style(r, 11.5, BODY)


def accent_rule(slide, x=0.5, y=1.33, w=1.4, h=0.055):
    """A short indigo rule under the slide title for visual hierarchy."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = HEAD
    sh.line.fill.background(); sh.shadow.inherit = False


def body_shape(slide):
    cand = [s for s in slide.shapes if s.has_text_frame and Emu(s.top).inches > 1.2
            and Emu(s.height).inches < 5.0]
    return max(cand, key=lambda s: Emu(s.width).inches * Emu(s.height).inches, default=None)


def content_panel(slide, x=0.45, y=1.5, w=9.1, h=3.52):
    """A light rounded content card with a left indigo accent stripe (the 'designed' look)."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = BORDER; panel.line.width = Pt(0.75); panel.shadow.inherit = False
    stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y + 0.16), Inches(0.09), Inches(h - 0.32))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = HEAD
    stripe.line.fill.background(); stripe.shadow.inherit = False
    return panel


def footer(slide, num, total=11):
    lb = slide.shapes.add_textbox(Inches(0.5), Inches(5.13), Inches(6.5), Inches(0.26))
    r0 = lb.text_frame.paragraphs[0].add_run()
    r0.text = "Team Argmax · Intelligent Candidate Discovery & Ranking"; style(r0, 8, MUTE)
    rb = slide.shapes.add_textbox(Inches(8.1), Inches(5.13), Inches(1.4), Inches(0.26))
    pr = rb.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.RIGHT
    rr = pr.add_run(); rr.text = f"{num:02d} / {total}"; style(rr, 8, MUTE)


def console_box(slide, x, y, w, h, lines):
    """A dark 'terminal' card showing real command output (evidence)."""
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bx.fill.solid(); bx.fill.fore_color.rgb = RGBColor(0x0F, 0x14, 0x20)
    bx.line.color.rgb = RGBColor(0x2A, 0x33, 0x45); bx.line.width = Pt(0.75); bx.shadow.inherit = False
    tf = bx.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    first = True
    for kind, text in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        no_bullet(p, marL=0.0); p.line_spacing = 1.12; p.space_after = Pt(2)
        r = p.add_run(); r.text = text
        r.font.name = "Consolas"; r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(0x7E, 0xE7, 0x9A) if kind == "cmd" else RGBColor(0xCE, 0xD6, 0xE3)
    return bx


def image_panel(slide, path, x, y, w, h, caption=None, cap_color=HEAD):
    """Place an image fit to (w,h) box, top-aligned and centered, with a caption + border."""
    from PIL import Image as _I
    iw, ih = _I.open(path).size
    ar = iw / ih
    draw_w = h * ar
    if draw_w > w:
        draw_w, draw_h = w, w / ar
    else:
        draw_h = h
    cx = x + (w - draw_w) / 2
    if caption:
        cap = slide.shapes.add_textbox(Inches(x), Inches(y - 0.32), Inches(w), Inches(0.3))
        rc = cap.text_frame.paragraphs[0]; rc.alignment = PP_ALIGN.CENTER
        rr = rc.add_run(); rr.text = caption; style(rr, 10, cap_color, bold=True)
    pic = slide.shapes.add_picture(path, Inches(cx), Inches(y), height=Inches(draw_h))
    pic.line.color.rgb = BORDER; pic.line.width = Pt(0.75)
    return pic


def stat_card(slide, x, y, w, h, number, lab):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = CARD
    c.line.color.rgb = BORDER; c.line.width = Pt(0.75); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = number; style(r1, 25, HEAD, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.0
    r2 = p2.add_run(); r2.text = lab; style(r2, 9.5, MUTE)


# ---- architecture diagram (native pptx shapes; renders crisply in the PDF) -------------------------
OFFLINE = RGBColor(0x0B, 0x6B, 0x66)   # teal
SCORED = RGBColor(0x1F, 0x4E, 0x79)    # blue
CACHE = RGBColor(0x92, 0x40, 0x0E)     # amber
GREYF = RGBColor(0xE9, 0xEC, 0xEF)     # light grey fill
ARROW = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def box(slide, x, y, w, h, text, fill, fg=WHITE, size=8):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = fg; r.font.bold = (i == 0)
    return sh


def arrow(slide, x1, y1, x2, y2, color=ARROW, dashed=False, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def label(slide, x, y, w, text, color, size=9):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    return tb


def draw_arch(slide):
    # Row A - offline precompute
    label(slide, 0.5, 1.58, 6.0, "1 · OFFLINE PRECOMPUTE  (GPU, one-time)", OFFLINE)
    a = ["100k\nprofiles", "Recall funnel\n(structured score)", "Bi-encoder\nbge-small\n(embeddings)",
         "Cross-encoder\nbge-reranker-v2-m3", "artifacts\n*.json\n(cache)"]
    ay, ah, aw, agap, ax0 = 1.92, 0.74, 1.55, 0.32, 0.5
    aboxes = []
    for i, t in enumerate(a):
        x = ax0 + i * (aw + agap)
        aboxes.append(box(slide, x, ay, aw, ah, t, CACHE if i == len(a) - 1 else OFFLINE))
        if i:
            px = ax0 + (i - 1) * (aw + agap) + aw
            arrow(slide, px, ay + ah / 2, x, ay + ah / 2)

    # Row B - scored ranking (rank.py)
    label(slide, 0.35, 3.2, 6.5, "2 · SCORED RANKING - rank.py (CPU · offline · <=5 min)", SCORED)
    b = ["candidates\n.jsonl", "Honeypot\nfilter", "JD rubric + features\nsemantic + rerank",
         "x penalties\nx behavioral", "Top-100\n+ reasoning", "submission\n.csv / .xlsx"]
    by, bh, bw, bgap, bx0 = 3.5, 0.74, 1.4, 0.18, 0.35
    bboxes = []
    for i, t in enumerate(b):
        x = bx0 + i * (bw + bgap)
        bboxes.append(box(slide, x, by, bw, bh, t, SCORED))
        if i:
            px = bx0 + (i - 1) * (bw + bgap) + bw
            arrow(slide, px, by + bh / 2, x, by + bh / 2)

    # cache -> scoring (dashed): from artifacts box bottom to the rubric box top
    cx = ax0 + 4 * (aw + agap) + aw / 2
    tx = bx0 + 2 * (bw + bgap) + bw / 2
    arrow(slide, cx, ay + ah, tx, by, color=CACHE, dashed=True, width=1.25)
    label(slide, 5.7, 2.86, 3.0, "cached scores -> loaded by rank.py", CACHE, size=7.5)

    # demo + compliance strip
    box(slide, 0.5, 4.46, 9.0, 0.56,
        "Demo sandbox: app/ (FastAPI) + frontend/ (React + Tailwind) -> single Docker image (HF Space), "
        "reading the same src/.   Compliance: GPU / transformers live only in precompute; rank.py is "
        "stdlib · CPU · offline.", GREYF, fg=RGBColor(0x20, 0x25, 0x2B), size=8)


def draw_results(slide):
    stats = [("0", "honeypots in top-100"), ("~150s", "CPU · offline run"),
             ("83/100", "in the 5-9 yr band"), ("97/100", "show ranking work")]
    y, h, w, gap, x0 = 1.6, 1.18, 2.13, 0.21, 0.45
    for i, (num, lab) in enumerate(stats):
        stat_card(slide, x0 + i * (w + gap), y, w, h, num, lab)
    console_box(slide, 0.45, 3.04, 9.1, 1.5, [
        ("cmd", "$ just check        # rank.py -> submission.csv, then validate"),
        ("out", "Scored 100,000 candidates (65 honeypots filtered) in 146.2s; top score 1.00, cutoff 0.75"),
        ("out", "Wrote 100 ranked candidates to submission.csv     Submission is valid."),
        ("cmd", "$ uv run python scripts/quality_proxy.py submission.csv"),
        ("out", "top-100:  in-band 83/100   ranking-work 97/100   juniors 1/100   available 92/100"),
    ])
    note = slide.shapes.add_textbox(Inches(0.5), Inches(4.64), Inches(9.0), Inches(0.32))
    rn = note.text_frame.paragraphs[0].add_run()
    rn.text = ("Top-10: senior product-company retrieval/ranking/NLP engineers (Meta, Apple, Netflix, "
               "Zomato, Paytm, Razorpay). Ground truth hidden - validated by ablation + calibration (docs/09).")
    style(rn, 9, MUTE, italic=True)


def panel_body(slide, items, x=0.78, y=1.64, w=8.45, h=3.25):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    fill_body(tb.text_frame, items)
    return tb


def _maprow(slide, y, left, right, kind):
    """One row of the JD->mechanism mapping: tone stripe + left (JD) -> right (our encoding)."""
    tone = {"want": HEAD, "accept": GREEN, "reject": RGBColor(0xE1, 0x1D, 0x48)}[kind]
    h = 0.46
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9.0), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = BORDER; card.line.width = Pt(0.75); card.shadow.inherit = False
    stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y + 0.08), Inches(0.07), Inches(h - 0.16))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = tone; stripe.line.fill.background(); stripe.shadow.inherit = False
    lt = slide.shapes.add_textbox(Inches(0.72), Inches(y), Inches(3.95), Inches(h))
    lt.text_frame.word_wrap = True; lt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    lt.text_frame.margin_top = lt.text_frame.margin_bottom = Inches(0.02)
    rl = lt.text_frame.paragraphs[0].add_run(); rl.text = left; style(rl, 9.5, BODY)
    ar = slide.shapes.add_textbox(Inches(4.68), Inches(y), Inches(0.42), Inches(h))
    ar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = ar.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
    ra = ap.add_run(); ra.text = "->"; style(ra, 11, MUTE, bold=True)
    rt = slide.shapes.add_textbox(Inches(5.15), Inches(y), Inches(4.25), Inches(h))
    rt.text_frame.word_wrap = True; rt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    rt.text_frame.margin_top = rt.text_frame.margin_bottom = Inches(0.02)
    rr = rt.text_frame.paragraphs[0].add_run(); rr.text = right; style(rr, 9.5, BODY)


def draw_jd_mapping(slide):
    label(slide, 0.72, 1.58, 3.95, "WHAT THE JD WANTS / REJECTS", MUTE, size=9)
    label(slide, 5.15, 1.58, 4.25, "HOW OUR ENGINE ENCODES IT", MUTE, size=9)
    rows = [
        ("Embeddings retrieval, vector/hybrid search, ranking eval (NDCG/MAP), strong Python",
         "Concept lexicon over career descriptions + role-gated skill depth", "want"),
        ("5-9 yrs, product company, shipped an end-to-end ranking / search / recsys",
         "Experience band + product-vs-services signal + semantic match", "want"),
        ("Tier-5 who built a recsys but never writes 'RAG' / 'Pinecone'",
         "bge embeddings over career text catch the plain-language fit", "accept"),
        ("'Marketing Manager' with a perfect AI skill list",
         "Engineering-role gate; the skills list is excluded from relevance", "reject"),
        ("Services-only career, CV/speech without NLP/IR, title-chaser",
         "Soft penalties  x0.45 / x0.6 / x0.85", "reject"),
        ("Stale login, 5% recruiter response (not actually available)",
         "Bounded behavioral-availability multiplier (down-weight)", "reject"),
    ]
    y0, step = 1.92, 0.5
    for i, (l, r, k) in enumerate(rows):
        _maprow(slide, y0 + i * step, l, r, k)


def main():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # slide 1 - title fields
    for sh in slides[0].shapes:
        if not sh.has_text_frame:
            continue
        for prefix, val in TITLE_FIELDS.items():
            if sh.text_frame.text.strip().startswith(prefix):
                tf = sh.text_frame; tf.word_wrap = True; tf.clear()
                tf.auto_size = MSO_AUTO_SIZE.NONE   # stop shrink-to-fit so all 3 labels stay 12pt
                p = tf.paragraphs[0]; p.line_spacing = 1.2   # avoid wrapped-line overlap
                # all three labels at the SAME size (matches the template); value in regular weight
                rl = p.add_run(); rl.text = f"{prefix} : "; style(rl, 12, BODY, bold=True)
                rv = p.add_run(); rv.text = val; style(rv, 12, BODY, bold=False)

    # content slides (2-6, 8, 9, 10)
    for n, items in CONTENT.items():
        slide = slides[n - 1]
        bs = body_shape(slide)            # remove the template's prompt-question box
        if bs is not None:
            bs._element.getparent().remove(bs._element)
        accent_rule(slide)

        if n == 8:
            draw_results(slide)
        elif n == 3:
            draw_jd_mapping(slide)
        elif n == 5:
            content_panel(slide, x=0.45, y=1.52, w=5.35, h=3.42)
            panel_body(slide, items, x=0.78, y=1.64, w=4.75, h=3.25)
            image_panel(slide, "docs/images/ui_detail.png", 6.0, 1.84, 3.5, 3.1,
                        caption="Explainability in the demo UI")
        elif n == 10:
            # full-page live-demo hero + a compact assets line
            a = slide.shapes.add_textbox(Inches(0.4), Inches(1.46), Inches(9.2), Inches(0.3))
            pa = a.text_frame.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
            ra = pa.add_run()
            ra.text = ("GitHub (public)    ·    ranked output XLSX    ·    approach deck (PDF)    ·    "
                       "docs/00-09 + PROJECT_LOG")
            style(ra, 10, BODY)
            c = slide.shapes.add_textbox(Inches(0.4), Inches(1.78), Inches(9.2), Inches(0.3))
            pc = c.text_frame.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
            rc = pc.add_run()
            rc.text = "LIVE on Hugging Face Space    ·    tsathya98-redrob-candidate-ranker.hf.space"
            style(rc, 10.5, GREEN, bold=True)
            image_panel(slide, "docs/images/ui_full.png", 0.8, 2.06, 8.4, 2.98)
        else:
            content_panel(slide)
            panel_body(slide, items)
        footer(slide, n)

    # slide 7 - System Architecture (native diagram)
    accent_rule(slides[6])
    draw_arch(slides[6])
    footer(slides[6], 7)

    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
